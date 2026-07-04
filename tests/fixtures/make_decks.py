"""Generate fixture decks for quality-engine tests.

Five contract-mandated fixture types: chart-heavy, photo-heavy,
text-heavy, filler-heavy, chart-as-image (the false-rejection bug case).
Decks are built with python-pptx (native charts produce real OOXML
graphicData chart parts); images are synthesized with PIL.
"""
from __future__ import annotations

import io
import random
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


def _photo_bytes(width: int = 640, height: int = 480, seed: int = 0) -> io.BytesIO:
    """Photo-like image: smooth gradients + noise = many colors, no
    straight edges, no uniform background."""
    rng = random.Random(seed)
    img = Image.new("RGB", (width, height))
    px = img.load()
    for y in range(height):
        for x in range(width):
            px[x, y] = (
                (x * 255 // width + rng.randrange(60)) % 256,
                (y * 255 // height + rng.randrange(60)) % 256,
                ((x + y) * 255 // (width + height) + rng.randrange(60)) % 256,
            )
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=88)
    buf.seek(0)
    return buf


def _chart_image_bytes(width: int = 640, height: int = 480, seed: int = 0) -> io.BytesIO:
    """Chart-like image: white background, black axes, few flat colors,
    labeled bars -- what an exported/pasted Excel chart looks like."""
    rng = random.Random(seed)
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    left, bottom, top, right = 80, height - 60, 40, width - 40
    draw.line([(left, top), (left, bottom)], fill="black", width=3)       # y axis
    draw.line([(left, bottom), (right, bottom)], fill="black", width=3)   # x axis
    for gy in range(top, bottom, (bottom - top) // 5):                    # gridlines
        draw.line([(left, gy), (right, gy)], fill=(200, 200, 200), width=1)
    colors = [(31, 119, 180), (255, 127, 14), (44, 160, 44)]
    n_bars = 6
    bar_w = (right - left) // (n_bars * 2)
    for i in range(n_bars):
        x0 = left + 20 + i * 2 * bar_w
        h = rng.randrange((bottom - top) // 3, bottom - top - 10)
        draw.rectangle([x0, bottom - h, x0 + bar_w, bottom],
                       fill=colors[i % 3], outline="black")
        draw.text((x0, bottom + 8), f"Q{i+1}", fill="black")
    draw.text((left, 10), "Revenue by quarter ($M)", fill="black")
    for i, v in enumerate(range(0, 120, 20)):
        draw.text((left - 40, bottom - i * (bottom - top) // 5 - 6), str(v), fill="black")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for bullet in bullets[1:]:
        para = body.add_paragraph()
        para.text = bullet
        para.level = 1


def _add_chart_slide(prs: Presentation, title: str, seed: int = 0) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rng = random.Random(seed)
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3", "Q4"]
    data.add_series("Revenue", tuple(rng.randrange(10, 100) for _ in range(4)))
    data.add_series("Cost", tuple(rng.randrange(5, 60) for _ in range(4)))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                           Inches(1), Inches(1.5), Inches(8), Inches(5), data)


def _add_image_slide(prs: Presentation, title: str, image_buf: io.BytesIO,
                     caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_buf, Inches(1), Inches(1.5), Inches(8), Inches(5))
    if caption:
        box = slide.shapes.add_textbox(Inches(1), Inches(6.6), Inches(8), Inches(0.5))
        box.text_frame.text = caption
        box.text_frame.paragraphs[0].font.size = Pt(12)


_LOREM = ("The strategic initiative delivered measurable operational improvements "
          "across all business units during the reporting period under review.")


def make_chart_heavy(path: Path) -> Path:
    prs = Presentation()
    _add_title_slide(prs, "FY2025 Results", "Investor Presentation")
    for i in range(5):
        _add_chart_slide(prs, f"Segment performance {i+1}", seed=i)
    for i in range(3):
        _add_text_slide(prs, f"Strategy pillar {i+1}", [_LOREM] * 4)
    _add_text_slide(prs, "Thank You", ["Questions?"])
    prs.save(path)
    return path


def make_photo_heavy(path: Path) -> Path:
    prs = Presentation()
    _add_title_slide(prs, "Brand Book", "Imagery")
    for i in range(7):
        _add_image_slide(prs, f"Moment {i+1}", _photo_bytes(seed=i))
    _add_text_slide(prs, "Thank You", ["-"])
    prs.save(path)
    return path


def make_text_heavy(path: Path) -> Path:
    prs = Presentation()
    _add_title_slide(prs, "Policy Handbook", "All-text edition")
    for i in range(8):
        _add_text_slide(prs, f"Policy Area {i+1}", [_LOREM] * 6)
    prs.save(path)
    return path


def make_filler_heavy(path: Path) -> Path:
    prs = Presentation()
    _add_title_slide(prs, "Quarterly Review", "Q4 2025")
    _add_text_slide(prs, "Agenda", ["Agenda"])
    slide = prs.slides.add_slide(prs.slide_layouts[2])   # section divider
    slide.shapes.title.text = "Section 1"
    _add_chart_slide(prs, "Key results", seed=42)
    _add_text_slide(prs, "Thank You", ["Thank you"])
    prs.save(path)
    return path


def make_chart_as_image(path: Path) -> Path:
    """THE regression fixture: charts pasted as raster images, zero native
    chart parts. The buggy prior generation rejected these as photo decks."""
    prs = Presentation()
    _add_title_slide(prs, "Market Analysis 2025", "Data Annex")
    for i in range(5):
        _add_image_slide(prs, f"Exhibit {i+1}",
                         _chart_image_bytes(seed=i),
                         caption=f"Figure {i+1}: revenue trends by region")
    for i in range(3):
        _add_text_slide(prs, f"Findings {i+1}", [_LOREM] * 3)
    _add_text_slide(prs, "Thank You", ["Q&A"])
    prs.save(path)
    return path


def make_all(out_dir: Path) -> dict[str, Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    return {
        "chart_heavy": make_chart_heavy(out_dir / "chart_heavy.pptx"),
        "photo_heavy": make_photo_heavy(out_dir / "photo_heavy.pptx"),
        "text_heavy": make_text_heavy(out_dir / "text_heavy.pptx"),
        "filler_heavy": make_filler_heavy(out_dir / "filler_heavy.pptx"),
        "chart_as_image": make_chart_as_image(out_dir / "chart_as_image.pptx"),
    }


if __name__ == "__main__":
    import sys
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "decks"
    for name, p in make_all(target).items():
        print(f"{name}: {p}")
