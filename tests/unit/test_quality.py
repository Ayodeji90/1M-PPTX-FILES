"""Quality engine tests against the five contract fixture decks."""
from __future__ import annotations

import json

from pptxsweeper.quality import classify, reclassify_from_vectors
from pptxsweeper.quality.images import classify_image


def test_chart_heavy_is_high_deliver(decks):
    report = classify(decks["chart_heavy"])
    assert report.error == ""
    assert report.slide_count == 10
    assert report.quality == "HIGH", report.explanations
    assert report.decision == "DELIVER", report.explanations
    assert report.chart_diagram_pages >= 3


def test_photo_heavy_is_rejected(decks):
    report = classify(decks["photo_heavy"])
    assert report.quality == "LOW", report.explanations
    assert report.decision == "REJECT"


def test_text_heavy_is_rejected(decks):
    report = classify(decks["text_heavy"])
    assert report.quality == "LOW", report.explanations
    assert report.decision == "REJECT"
    assert report.text_only_pct >= 0.75


def test_filler_excluded_from_denominator(decks):
    report = classify(decks["filler_heavy"])
    # 5 slides, only 1 substantive (the chart slide): title/agenda/
    # divider/thank-you must all be excluded from the denominator.
    assert report.slide_count == 5
    assert report.content_slide_count <= 2, [
        (s.index, s.is_structural_filler, s.filler_reason) for s in report.slides]
    assert report.quality in ("HIGH", "MEDIUM"), report.explanations
    assert report.decision in ("DELIVER", "REVIEW")


def test_chart_as_image_not_falsely_rejected(decks):
    """THE regression test for the prior generation's bug: decks whose
    charts are pasted raster images must not be classified photo-heavy."""
    report = classify(decks["chart_as_image"])
    assert report.quality in ("HIGH", "MEDIUM"), report.explanations
    assert report.decision in ("DELIVER", "REVIEW"), report.explanations
    assert report.photo_heavy_pct < 0.5, report.explanations
    analytical_images = sum(s.image_analytical_count for s in report.slides)
    assert analytical_images >= 3, "chart-like images must classify analytical"


def test_image_classifier_chart_vs_photo():
    from fixtures.make_decks import _chart_image_bytes, _photo_bytes
    chart_sig = classify_image(_chart_image_bytes(seed=7).getvalue())
    assert chart_sig.label == "analytical", chart_sig
    photo_sig = classify_image(_photo_bytes(seed=7).getvalue())
    assert photo_sig.label != "analytical", photo_sig


def test_ambiguous_image_never_counts_as_photo():
    # A tiny image is undecodable/noise -> ambiguous, not photo.
    sig = classify_image(b"not an image at all")
    assert sig.label == "ambiguous"


def test_min_slides_rejected(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"S{i}"
    p = tmp_path / "short.pptx"
    prs.save(p)
    report = classify(p)
    assert report.decision == "REJECT"
    assert "below contractual minimum" in report.explanations[0]


def test_reclassify_from_stored_vectors_matches(decks):
    """Threshold re-runs must work from persisted vectors alone."""
    report = classify(decks["chart_heavy"])
    stored = json.loads(json.dumps(report.feature_vectors_json()))
    replay = reclassify_from_vectors(stored, fmt="pptx")
    assert replay.quality == report.quality
    assert replay.decision == report.decision

    # tightened thresholds change the outcome without touching the file
    strict = reclassify_from_vectors(
        stored, {"high": {"min_chart_diagram_pages": 99},
                 "medium": {"min_chart_diagram_pages": 99}})
    assert strict.quality == "LOW"


def test_borderline_goes_to_review():
    from pptxsweeper.quality.report import SlideFeatures
    from pptxsweeper.quality.decide import decide_from_features
    # exactly 50% analytical == HIGH threshold -> within epsilon -> REVIEW
    slides = []
    for i in range(10):
        f = SlideFeatures(index=i, text_char_count=500)
        if i < 5:
            f.native_chart_count = 1
        slides.append(f)
    report = decide_from_features(slides)
    assert report.borderline is True
    assert report.decision == "REVIEW", report.explanations


def test_corrupt_pptx_rejects_not_raises(tmp_path):
    bad = tmp_path / "corrupt.pptx"
    bad.write_bytes(b"PK\x03\x04" + b"\x00" * 100)
    report = classify(bad)
    assert report.decision == "REJECT"
    assert report.error
