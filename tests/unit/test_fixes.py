"""Regression tests for the CPU/memory/contract fixes:

- composition contract: no ZeroDivisionError on degenerate config, and the
  70/30 HIGH/MEDIUM floor is enforced by the batch selector
- PPTXSWEEPER_OVERRIDE deep-merge of a VM config profile
- speaker notes reach the compliance text (PII screening) without touching
  the quality verdict
- docProps (core/app) retained with ONE zip open
- raster-image classification skipped on slides that already carry a
  native analytical object (the big classify CPU saving)
"""
from __future__ import annotations

import pytest


def _cand(i: int, quality: str, delivered: bool = False) -> dict:
    return {"id": i, "quality": quality,
            "delivered_filename": f"BATCH_01_file_{i:05d}.pptx" if delivered else None}


# --- composition contract -------------------------------------------------
def test_select_for_batch_no_crash_on_degenerate_config():
    from pptxsweeper.packager.compose import select_for_batch
    # The old shipped config was high_min_pct=0.0 / medium_max_pct=1.0,
    # which divided by zero and crashed the full-batch packager.
    sel = select_for_batch([], 100, high_min_pct=0.0, medium_max_pct=1.0)
    assert sel.count == 0
    # Degenerate values must not poison composition_ok either.
    rows = [_cand(i, "HIGH") for i in range(10)] + [_cand(i, "MEDIUM") for i in range(10, 20)]
    sel = select_for_batch(rows, 20, high_min_pct=0.0, medium_max_pct=1.0)
    assert sel.count == 20
    assert sel.composition_ok(0.0, 1.0) is True


def test_select_for_batch_enforces_70_30_floor():
    from pptxsweeper.packager.compose import select_for_batch
    rows = [_cand(i, "HIGH") for i in range(10)] + [_cand(i, "MEDIUM") for i in range(10, 20)]
    sel = select_for_batch(rows, 20, high_min_pct=0.70, medium_max_pct=0.30)
    # MEDIUM <= floor(high * 0.3/0.7) = floor(10*0.4286) = 4; 6 held back.
    assert len(sel.high) == 10
    assert len(sel.medium) == 4
    assert len(sel.surplus_medium) == 6
    assert sel.composition_ok(0.70, 0.30) is True
    # a short batch keeps the floor too: 3 HIGH + 20 MEDIUM -> only
    # floor(3*0.4286)=1 MEDIUM joins, rest held in reserve
    rows = [_cand(i, "HIGH") for i in range(3)] + [_cand(i, "MEDIUM") for i in range(3, 23)]
    sel = select_for_batch(rows, 25, high_min_pct=0.70, medium_max_pct=0.30)
    assert len(sel.high) == 3 and len(sel.medium) == 1
    assert len(sel.surplus_medium) == 19


def test_resolve_composition_rejects_degenerate_config():
    from pptxsweeper.packager.compose import resolve_composition
    # the old broken shipped config -> enforced contract defaults
    assert resolve_composition({"high_min_pct": 0.0, "medium_max_pct": 1.0}) \
        == (0.70, 0.20, 0.30)
    # anything WEAKER than the contract is clamped back to it
    assert resolve_composition({"high_min_pct": 0.60, "medium_max_pct": 0.35}) \
        == (0.70, 0.20, 0.30)
    # a sane config passes through untouched
    assert resolve_composition({"high_min_pct": 0.70, "medium_min_pct": 0.20,
                                "medium_max_pct": 0.30}) == (0.70, 0.20, 0.30)
    # stricter-than-contract passes through
    assert resolve_composition({"high_min_pct": 0.80, "medium_max_pct": 0.25}) \
        == (0.80, 0.20, 0.25)


def test_package_stage_composition_resolves(decks, registry):
    """PackageStage must come up with contract values even if config.yaml
    is later hand-edited back to a degenerate state."""
    from pptxsweeper.packager.compose import resolve_composition
    from pptxsweeper.packager.package_stage import PackageStage
    from pptxsweeper.config import Config
    cfg = Config.load()
    cfg.raw["batch"]["composition"]["high_min_pct"] = 0.0  # simulate bad edit
    stage = PackageStage(cfg, registry)
    assert stage.high_min_pct == 0.70
    assert stage.medium_max_pct == 0.30


# --- speaker notes reach compliance text ---------------------------------
def test_speaker_notes_in_full_text(tmp_path):
    from pptx import Presentation
    from pptxsweeper.quality import classify

    # Blank slides: the ONLY text in the deck lives in the speaker notes,
    # so we can prove notes reach the compliance text while never counting
    # as slide text in the quality metrics.
    prs = Presentation()
    for i in range(6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        slide.notes_slide.notes_text_frame.text = \
            "internal only: pii@example.com, john.doe@gmail.com"
    p = tmp_path / "notes.pptx"
    prs.save(p)

    report = classify(p)
    # notes text reaches the compliance text (PII screens see it)...
    assert "pii@example.com" in report.full_text
    assert "john.doe@gmail.com" in report.full_text
    # ...but never counts as slide text (quality math is body-text-only)
    assert all(s.text_char_count == 0 for s in report.slides)


# --- docProps retained (single zip open) ---------------------------------
def test_doc_properties_retained(tmp_path):
    from pptx import Presentation
    from pptxsweeper.quality import classify

    prs = Presentation()
    for i in range(6):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {i}"
    prs.core_properties.title = "Retained Title"
    prs.core_properties.author = "Test Author"
    p = tmp_path / "meta.pptx"
    prs.save(p)

    report = classify(p)
    assert report.doc_properties.get("title") == "Retained Title"
    assert report.doc_properties.get("creator") == "Test Author"


# --- image classification skipped on native-analytical slides ------------
def test_image_classification_skipped_on_native_analytical_slide(tmp_path):
    from pptx import Presentation
    from fixtures.make_decks import _photo_bytes
    from pptxsweeper.quality.ooxml import parse_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Quarterly results"
    slide.shapes.add_table(rows=3, cols=3, left=1_000_000, top=1_000_000,
                           width=5_000_000, height=2_000_000)   # native analytical
    slide.shapes.add_picture(_photo_bytes(seed=3),               # photo on same slide
                             1_000_000, 3_500_000)
    p = tmp_path / "table_photo.pptx"
    prs.save(p)

    def boom(_data):
        raise AssertionError("image classifier must NOT run on native-analytical slides")

    features, _texts, signals, _notes, _props = parse_pptx(str(p), boom)
    assert features[0].table_count == 1
    assert features[0].image_count == 1          # image still counted
    assert features[0].image_photo_count == 0    # but never penalized as photo
    labels = {s["label"] for s in signals[0]}
    assert labels == {"unclassified"}


def test_chart_as_image_still_classified(decks):
    """The chart-as-image deck has NO native objects: its raster charts must
    still run through the classifier (regression guard for the skip)."""
    from pptxsweeper.quality.ooxml import parse_pptx
    features, _t, signals, _n, _p = parse_pptx(str(decks["chart_as_image"]))
    analytical = sum(1 for slide_sigs in signals
                     for s in slide_sigs if s.get("label") == "analytical")
    assert analytical >= 3
    assert all(s.get("label") != "unclassified" for slide_sigs in signals
               for s in slide_sigs)


# --- config overlay deep-merge --------------------------------------------
def test_config_override_deep_merge(tmp_path, monkeypatch):
    from pptxsweeper.config import Config
    (tmp_path / "config.yaml").write_text("a:\n  x: 1\n  y: 2\nb: 3\n")
    (tmp_path / "override.yaml").write_text("a:\n  y: 99\n")
    monkeypatch.setenv("PPTXSWEEPER_OVERRIDE", "override.yaml")
    cfg = Config.load(start=tmp_path)
    assert cfg.raw["a"]["x"] == 1      # untouched base key
    assert cfg.raw["a"]["y"] == 99     # overridden
    assert cfg.raw["b"] == 3           # untouched base section


def test_config_override_from_dotenv(tmp_path, monkeypatch):
    """PPTXSWEEPER_OVERRIDE set in .env (the documented VM mechanism)
    must be honored -- not just real shell env vars."""
    from pptxsweeper.config import Config
    (tmp_path / "config.yaml").write_text("a:\n  x: 1\n  y: 2\nb: 3\n")
    (tmp_path / "override.yaml").write_text("a:\n  y: 99\n")
    (tmp_path / ".env").write_text("PPTXSWEEPER_OVERRIDE=override.yaml\n")
    monkeypatch.delenv("PPTXSWEEPER_OVERRIDE", raising=False)
    cfg = Config.load(start=tmp_path)
    assert cfg.raw["a"]["x"] == 1
    assert cfg.raw["a"]["y"] == 99
    assert cfg.raw["b"] == 3


def test_config_override_missing_raises(tmp_path, monkeypatch):
    from pptxsweeper.config import Config, ConfigError
    (tmp_path / "config.yaml").write_text("a: 1\n")
    monkeypatch.setenv("PPTXSWEEPER_OVERRIDE", "nope.yaml")
    with pytest.raises(ConfigError):
        Config.load(start=tmp_path)


# --- harvester / filter fixes (found live on the VM) ----------------------
def test_govdata_relative_resource_urls_absolutized():
    """open.canada.ca returns path-only resource URLs; they must be
    absolutized against the portal base (and un-joinable URLs dropped)."""
    from pptxsweeper.harvesters.tier7_govdata import _absolute_resource_url
    portal = "https://open.canada.ca/data/en"
    assert _absolute_resource_url(portal, "/data/dataset/abc/download/d.pptx") \
        == "https://open.canada.ca/data/dataset/abc/download/d.pptx"
    # absolute URLs pass through untouched
    assert _absolute_resource_url(portal, "https://cdn.example.com/a.pptx") \
        == "https://cdn.example.com/a.pptx"
    # empty / hopeless inputs are dropped
    assert _absolute_resource_url(portal, "") == ""
    assert _absolute_resource_url(portal, "javascript:void(0)") == ""


def test_filter_rejects_malformed_urls(tmp_path, registry):
    """Scheme-less / host-less URLs must be filtered out in ONE pass instead
    of burning download attempts."""
    from pptxsweeper.config import Config
    from pptxsweeper.stages.filter_stage import run_filter

    (tmp_path / "block.txt").write_text("")
    raw = {
        "compliance": {"excluded_sources_file": None,
                        "blocklist_file": "block.txt"},
        "allowed_formats": ["pptx", "ppt"],
        "filter": {"per_domain_cap": 10000, "tier_domain_caps": {}},
    }
    cfg = Config(raw, root=tmp_path, env_path=None)
    reg = registry
    reg.conn.executemany(
        "INSERT INTO urls (url, domain, tier, discovery_source, status) "
        "VALUES (?,?,?,?,?)",
        [
            ("/data/dataset/abc/download/x.pptx", "", 7, "govdata:open.canada.ca", "discovered"),
            ("ftp://files.example.com/a.pptx", "files.example.com", 7, "govdata:data.gov.uk", "discovered"),
            ("https://example.com/good.pptx", "example.com", 1, "wayback_cdx", "discovered"),
        ],
    )
    stats = run_filter(cfg, reg)
    assert stats["malformed_url"] == 2
    statuses = {r["url"]: r["status"] for r in
                reg.conn.execute("SELECT url, status FROM urls").fetchall()}
    assert statuses["/data/dataset/abc/download/x.pptx"] == "filtered_out"
    assert statuses["ftp://files.example.com/a.pptx"] == "filtered_out"
    assert statuses["https://example.com/good.pptx"] == "discovered"


def test_ir_harvester_none_website_does_not_crash(tmp_path):
    """A seed row whose `website` field is None must not crash discover()
    (regression for the .strip() AttributeError seen on the VM)."""
    import asyncio
    from pptxsweeper.harvesters.tier2_ir import InvestorRelationsHarvester

    # NOTE: a SHORT row (missing trailing field) is what makes csv.DictReader
    # yield website=None -- a literal "None" cell would be the truthy string
    # "None" and would never reproduce the original crash.
    seeds = tmp_path / "seeds.csv"
    seeds.write_text("ticker,company,website\nAAA,Company A,https://a.example\nBBB,Company B\n")

    h = InvestorRelationsHarvester.__new__(InvestorRelationsHarvester)
    h.cfg = type("Cfg", (), {"raw": {"harvesters": {"tier2": {"seed_tickers_file": str(seeds)}}},
                            "root": tmp_path})()
    h.owns_domain = lambda _d: False   # never actually fetch

    async def _go():
        n = 0
        async for _ in h.discover():
            n += 1
        return n

    assert asyncio.run(_go()) == 0   # no crash, no candidates


# --- wayback streaming size gates -----------------------------------------
def test_wayback_oversize_returns_rejected_signal(tmp_path):
    """fetch_to_file must stream (not buffer), enforce max_bytes, and delete
    the partial file -- the RAM-spike fix for the fallback path."""
    import asyncio
    import httpx

    from pptxsweeper.net.wayback import WaybackFetcher

    async def _run():
        # httpx MockTransport that streams a body larger than max_bytes
        body = b"x" * 1000

        def handler(request: httpx.Request) -> httpx.Response:
            return httpx.Response(200, content=body)

        client = httpx.AsyncClient(transport=httpx.MockTransport(handler))
        wb = WaybackFetcher(client, fetch_base_url="https://web.archive.org",
                            requests_per_sec=1000.0)
        dest = tmp_path / "wb.part"
        ok, sha, size, snap, status = await wb.fetch_to_file(
            "http://example.com/a.ppt", dest, timestamp="20250101000000",
            max_bytes=100, min_bytes=1)
        await client.aclose()
        return ok, sha, size, snap, status, dest.exists()

    ok, sha, size, snap, status, exists = asyncio.run(_run())
    assert ok is False
    assert size == 1000          # size reported for the reject reason
    assert exists is False       # partial file deleted
